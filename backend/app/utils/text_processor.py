import os
from typing import List
import PyPDF2
from pptx import Presentation as PPTXPresentation

def extract_text_from_presentation(file_path: str) -> str:
    """
    Extract text from a presentation file (PDF or PowerPoint).
    
    Args:
        file_path: Path to the presentation file
        
    Returns:
        str: Extracted text from the presentation
        
    Raises:
        Exception: If file type is not supported or extraction fails
    """
    file_ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_ext == '.pdf':
            return _extract_text_from_pdf(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return _extract_text_from_pptx(file_path)
        else:
            raise Exception(f"Unsupported file type: {file_ext}")
    except Exception as e:
        raise Exception(f"Failed to extract text: {str(e)}")

def _extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file."""
    text = []
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text.append(page.extract_text())
    return "\n\n".join(text)

def _extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint file."""
    text = []
    prs = PPTXPresentation(file_path)
    
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        text.append("\n".join(slide_text))
    
    return "\n\n".join(text)

def chunk_text(text: str, max_chunk_size: int = 1000, overlap: int = 100) -> List[str]:
    """
    Split text into overlapping chunks of approximately equal size.
    
    Args:
        text: Text to split into chunks
        max_chunk_size: Maximum size of each chunk in characters
        overlap: Number of characters to overlap between chunks
        
    Returns:
        List[str]: List of text chunks
    """
    # Split text into sentences (simple approach)
    sentences = text.replace('\n', ' ').split('. ')
    sentences = [s.strip() + '.' for s in sentences if s.strip()]
    
    chunks = []
    current_chunk = []
    current_size = 0
    
    for sentence in sentences:
        sentence_size = len(sentence)
        
        # If adding this sentence would exceed max size, save current chunk
        if current_size + sentence_size > max_chunk_size and current_chunk:
            chunks.append(' '.join(current_chunk))
            
            # Keep last few sentences for overlap
            overlap_size = 0
            overlap_sentences = []
            for s in reversed(current_chunk):
                if overlap_size + len(s) > overlap:
                    break
                overlap_sentences.insert(0, s)
                overlap_size += len(s)
            
            current_chunk = overlap_sentences
            current_size = overlap_size
        
        current_chunk.append(sentence)
        current_size += sentence_size
    
    # Add the last chunk if it exists
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks 